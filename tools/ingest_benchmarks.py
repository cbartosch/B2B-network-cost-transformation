#!/usr/bin/env python3
"""Ingest a folder of benchmark sources, whatever format they are in.

Conversion happens here, on your machine, and only the extracted text is sent
to the API. A confidential source file therefore never has to enter the
container or the database - only the text and the structured observations,
which carry a rights flag. That is also why this is a script rather than an
upload endpoint: the API image has no pptx, xlsx or pdf parser and adding one
would move client files into a service that has no reason to hold them.

Usage:

    python tools/ingest_benchmarks.py ./benchmarks --rights PRIOR_ENGAGEMENT --org "<prior engagement>"
    python tools/ingest_benchmarks.py ./public-tariffs --rights PUBLISHED
    python tools/ingest_benchmarks.py ./one.pptx --dry-run      # show the text only

Supported: .pptx .docx .xlsx .csv .pdf .md .txt .html

Every file is converted with whatever library is present; missing converters
are reported per file rather than aborting the run, so a folder of mixed
formats ingests what it can and tells you what it could not. Install what you
need locally:

    pip install python-pptx python-docx openpyxl pypdf

RIGHTS. --rights is mandatory and there is no default, because the difference
between a published tariff and a prior engagement's RFP responses is the
difference between a benchmark and another client's commercial position.
PRIOR_ENGAGEMENT observations land uncleared and contribute to nothing until
someone clears them by name.
"""
import argparse
import json
import pathlib
import sys
import urllib.error
import urllib.request

TEXT_SUFFIXES = {".md", ".txt", ".csv", ".html", ".htm", ".json"}


def _pptx(path):
    from pptx import Presentation
    out = []
    for i, slide in enumerate(Presentation(path).slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
        if parts:
            out.append(f"--- slide {i} ---\n" + "\n".join(parts))
    return "\n\n".join(out)


def _docx(path):
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for t in d.tables:
        for row in t.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    out = []
    for ws in wb.worksheets:
        rows = [" | ".join("" if c is None else str(c) for c in r)
                for r in ws.iter_rows(values_only=True)]
        rows = [r for r in rows if r.strip(" |")]
        if rows:
            out.append(f"--- sheet {ws.title} ---\n" + "\n".join(rows))
    return "\n\n".join(out)


def _pdf(path):
    from pypdf import PdfReader
    return "\n\n".join(
        f"--- page {i} ---\n{(pg.extract_text() or '')}"
        for i, pg in enumerate(PdfReader(path).pages, start=1))


CONVERTERS = {".pptx": _pptx, ".docx": _docx, ".xlsx": _xlsx, ".xlsm": _xlsx,
              ".pdf": _pdf}


def to_text(path: pathlib.Path) -> str:
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return path.read_text(encoding="utf-8", errors="replace")
    fn = CONVERTERS.get(suffix)
    if fn is None:
        raise RuntimeError(f"no converter for {suffix}")
    return fn(str(path))


def post(base: str, token: str | None, payload: dict, timeout: float) -> dict:
    req = urllib.request.Request(
        f"{base}/v1/benchmarks:extract",
        data=json.dumps(payload).encode(),
        headers={"Content-Type": "application/json",
                 **({"X-API-Token": token} if token else {})})
    try:
        with urllib.request.urlopen(req, timeout=timeout) as r:
            return json.load(r)
    except urllib.error.HTTPError as e:
        return {"_error": f"HTTP {e.code}: {e.read().decode()[:600]}"}
    except urllib.error.URLError as e:
        return {"_error": f"unreachable: {e.reason}"}


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("path", help="a file, or a folder to walk")
    ap.add_argument("--rights", required=True,
                    choices=["PUBLISHED", "VENDOR_SUPPLIED", "PRIOR_ENGAGEMENT"],
                    help="how this material may be used. No default: a "
                         "published tariff and a prior engagement's RFP "
                         "responses are not the same thing")
    ap.add_argument("--org", help="whose material this is")
    ap.add_argument("--as-of", help="vintage, e.g. 2026-05-14")
    ap.add_argument("--api", default="http://localhost:8000")
    ap.add_argument("--token", help="API_TOKEN, if the API enforces one")
    ap.add_argument("--timeout", type=float, default=600.0)
    ap.add_argument("--max-chars", type=int, default=40000,
                    help="split larger sources into chunks; an over-long "
                         "source truncates the model's reply mid-list")
    ap.add_argument("--dry-run", action="store_true",
                    help="print the extracted text and send nothing")
    args = ap.parse_args()

    root = pathlib.Path(args.path).expanduser()
    files = ([root] if root.is_file()
             else sorted(p for p in root.rglob("*") if p.is_file()))
    if not files:
        print(f"no files under {root}")
        return 1

    ok = failed = skipped = 0
    for path in files:
        try:
            text = to_text(path).strip()
        except RuntimeError as exc:
            print(f"SKIP  {path.name}: {exc}")
            skipped += 1
            continue
        except Exception as exc:                              # noqa: BLE001
            print(f"FAIL  {path.name}: {type(exc).__name__}: {exc}")
            failed += 1
            continue
        if not text:
            print(f"SKIP  {path.name}: no extractable text")
            skipped += 1
            continue

        chunks = [text[i:i + args.max_chars]
                  for i in range(0, len(text), args.max_chars)]
        for n, chunk in enumerate(chunks, start=1):
            label = f"{path.name}" + (f" [{n}/{len(chunks)}]" if len(chunks) > 1 else "")
            if args.dry_run:
                print(f"\n===== {label} ({len(chunk)} chars) =====\n{chunk[:1500]}")
                continue
            res = post(args.api, args.token, {
                "text": chunk, "source_document": path.name,
                "source_locator": f"chunk {n}/{len(chunks)}" if len(chunks) > 1 else None,
                "source_org": args.org, "rights_basis": args.rights,
                "as_of": args.as_of}, args.timeout)
            if "_error" in res:
                print(f"FAIL  {label}: {res['_error']}")
                failed += 1
            else:
                print(f"OK    {label}: {res['stored']} observation(s)"
                      + (f", {len(res['rejected'])} rejected" if res.get("rejected") else "")
                      + ("" if res.get("rights_cleared") else "  [rights NOT cleared]"))
                ok += 1

    if not args.dry_run:
        print(f"\n{ok} extracted, {failed} failed, {skipped} skipped")
        print("Nothing is priced yet. Review with:")
        print("  GET  /v1/benchmarks/observations?rights_cleared=false")
        print("  POST /v1/benchmarks/bands:derive   (dry_run defaults to true)")
    return 0 if not failed else 1


if __name__ == "__main__":
    sys.exit(main())
