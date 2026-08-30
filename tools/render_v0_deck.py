#!/usr/bin/env python3
"""Render a V0 estimate snapshot as a PowerPoint deck.

Runs locally, like tools/ingest_benchmarks.py and for the same reasons: the
API image carries no pptx library and has no business holding one, and a
deliverable belongs on the machine where the client template lives.

    python tools/render_v0_deck.py --case <case-id> -o v0.pptx
    python tools/render_v0_deck.py --case <case-id> --snapshot <id> -o v0.pptx

Everything on these slides is read from the stored snapshot - current TCO,
scenarios, per-lever savings, confidence with the ceilings that were applied,
coverage, and the provenance of each. Nothing is recomputed here, so the deck
cannot disagree with the estimate it reports.

**What this deliberately does not generate.** The DHL reference deck this was
modelled on has slides on strategic fit, target architecture and design
guardrails. Those are consulting judgement - a reading of Strategy 2030, a
view on route diversity, an architecture stance. The workbench holds no
evidence for any of them, and a generator that emitted them anyway would be
manufacturing confident narrative from a cost model, which is the one thing
this system refuses everywhere else. The deck ends with a slide stating what
the estimate rests on and what it is not, and the narrative slides remain the
consultant's to write.

Dependency: pip install python-pptx
"""
import argparse
import json
import pathlib
import sys
import urllib.error
import urllib.request
from datetime import date

INK = "1A1A1A"
MUTED = "6B6B6B"
ACCENT = "C8102E"
RULE = "D8D8D8"


def _get(base, path, token, timeout=60.0):
    req = urllib.request.Request(
        f"{base}{path}", headers={"X-API-Token": token} if token else {})
    try:
        with urllib.request.urlopen(req, timeout=timeout) as r:
            return json.load(r)
    except urllib.error.HTTPError as e:
        raise SystemExit(f"API returned HTTP {e.code}: {e.read().decode()[:400]}")
    except urllib.error.URLError as e:
        raise SystemExit(f"API unreachable: {e.reason}")


def _money(v, currency="EUR"):
    try:
        f = float(v)
    except (TypeError, ValueError):
        return str(v)
    for cut, suffix in ((1e9, "B"), (1e6, "M"), (1e3, "k")):
        if abs(f) >= cut:
            return f"{currency} {f / cut:,.1f}{suffix}"
    return f"{currency} {f:,.0f}"


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--case", required=True)
    ap.add_argument("--snapshot", help="defaults to the most recent")
    ap.add_argument("-o", "--out", default="v0_estimate.pptx")
    ap.add_argument("--api", default="http://localhost:8000")
    ap.add_argument("--token")
    args = ap.parse_args()

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Emu, Inches, Pt
    except ImportError:
        raise SystemExit("python-pptx is required: pip install python-pptx")

    case = _get(args.api, f"/v1/outside-in/cases/{args.case}", args.token)
    snaps = _get(args.api, f"/v1/outside-in/cases/{args.case}/estimates",
                 args.token).get("snapshots", [])
    if not snaps:
        raise SystemExit("no estimate snapshots on this case - run V0 first")
    snap = (next((x for x in snaps
                  if x["estimate_snapshot_id"] == args.snapshot), None)
            if args.snapshot else snaps[0])
    if snap is None:
        raise SystemExit(f"snapshot {args.snapshot} not found on this case")
    disp = _get(args.api, f"/v1/outside-in/cases/{args.case}/domain-dispositions",
                args.token)

    currency = (case.get("base_currency") or "EUR")
    pins = snap.get("pins") or {}
    method = pins.get("estimate_method", "BUILD_UP")
    cov = snap.get("coverage") or {}
    conf = snap.get("confidence") or {}
    scen = snap.get("scenarios") or {}

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide(kicker, title, standfirst=None):
        sl = prs.slides.add_slide(blank)
        if kicker:
            _text(sl, kicker.upper(), 0.7, 0.45, 11.9, 0.3, 11, MUTED, bold=True,
                  spacing=2)
        _text(sl, title, 0.7, 0.8, 11.9, 0.9, 26, INK, bold=True)
        if standfirst:
            _text(sl, standfirst, 0.7, 1.7, 11.9, 0.5, 13, MUTED)
        return sl

    def _text(sl, s, x, y, w, h, size, color, bold=False, spacing=None,
              align=None):
        box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = s
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)
        if spacing:
            run.font._rPr.set("spc", str(int(spacing * 100)))
        if align is not None:
            p.alignment = align
        return box

    def stat(sl, x, y, value, label, note=None, w=2.7):
        # A range like "EUR 31.0M-EUR 58.0M" is three times the width of a
        # single figure. At a fixed 30pt it wrapped onto the label beneath it -
        # the most common and most visible defect in a generated deck, so the
        # size follows the string rather than the other way round.
        text = str(value)
        size = 30 if len(text) <= 11 else (21 if len(text) <= 20 else 16)
        _text(sl, text, x, y, w, 0.75, size, ACCENT, bold=True)
        _text(sl, label, x, y + 0.78, w, 0.3, 11, INK, bold=True)
        if note:
            _text(sl, note, x, y + 1.08, w, 0.9, 9.5, MUTED)

    def table(sl, rows, x, y, w, col_w=None, size=10.5):
        n, m = len(rows), len(rows[0])
        shape = sl.shapes.add_table(n, m, Inches(x), Inches(y), Inches(w),
                                    Inches(0.32 * n))
        tbl = shape.table
        if col_w:
            for i, cw in enumerate(col_w):
                tbl.columns[i].width = Emu(int(Inches(cw)))
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.text = str(val)
                para = cell.text_frame.paragraphs[0]
                for run in para.runs:
                    run.font.size = Pt(size)
                    run.font.bold = (r == 0)
                    run.font.color.rgb = RGBColor.from_string(
                        INK if r == 0 else MUTED)
        return tbl

    # ---------------------------------------------------------------- 1 title
    sl = prs.slides.add_slide(blank)
    _text(sl, case.get("subject_entity_legal_name") or "Subject entity",
          0.7, 2.3, 11.9, 0.9, 32, INK, bold=True)
    _text(sl, f"Outside-in network cost estimate - V0 ({method})",
          0.7, 3.3, 11.9, 0.5, 16, MUTED)
    _text(sl, f"{date.today().isoformat()}   |   snapshot "
              f"{snap['estimate_snapshot_id'][:8]}   |   status "
              f"{snap.get('v0_status')}", 0.7, 6.6, 11.9, 0.4, 11, MUTED)

    # ------------------------------------------------------------ 2 headline
    headline = max(scen, key=lambda k: float(
        scen[k]["gross_run_rate_savings"]["base"])) if scen else None
    sl = slide("Executive view", "Modelled annual run-rate savings",
               "Every figure on this deck is read from the stored estimate "
               "snapshot. Nothing is recomputed here, so the deck cannot "
               "disagree with the estimate it reports.")
    if headline:
        g = scen[headline]["gross_run_rate_savings"]
        stat(sl, 0.7, 2.4, f"{currency} {float(g['low'])/1e6:,.0f}-"
                            f"{float(g['high'])/1e6:,.0f}M",
             "gross run-rate savings",
             f"Scenario {headline}. Low/high are the modelled range, not a "
             f"confidence interval.")
        stat(sl, 3.8, 2.4, _money(g["base"], currency), "base case",
             "Use as the planning figure.")
        stat(sl, 6.9, 2.4, _money((snap.get("current_tco") or {}).get("total"),
                                  currency), "modelled current cost",
             cov.get("spend_basis", "").replace("_", " ").lower() or None)
        stat(sl, 10.0, 2.4, str(conf.get("band", "-")), "confidence band",
             f"score {conf.get('score', '-')}; ceilings applied where the "
             f"evidence does not support more.")
    _text(sl, f"V0 status: {snap.get('v0_status')} - {cov.get('reason', '')}",
          0.7, 4.6, 11.9, 0.8, 11, MUTED)

    # --------------------------------------------------------------- 3 basis
    sl = slide("Basis", "What this estimate is built on")
    if method == "ANCHOR":
        basis = pins.get("anchor_basis") or {}
        pool = basis.get("addressable_pool") or {}
        share = basis.get("addressable_share") or {}
        stat(sl, 0.7, 2.4, _money(basis.get("anchor_value"), currency),
             "disclosed anchor", "An upper bound: carries voice, mobile and "
                                 "non-WAN services.")
        stat(sl, 3.8, 2.4,
             f"{float(share.get('low', 0)):.0%}-{float(share.get('high', 0)):.0%}",
             "addressable share", "A governed assumption, not an observation.")
        stat(sl, 6.9, 2.4,
             f"{currency} {float(pool.get('low', 0))/1e6:,.0f}-"
             f"{float(pool.get('high', 0))/1e6:,.0f}M",
             "addressable pool", "Anchor x share. The levers act on this.")
        stat(sl, 10.0, 2.4, str(basis.get("anchor_origin", "-")).replace("_", " "),
             "anchor provenance",
             "A typed anchor makes the whole estimate an assertion.")
        _text(sl, basis.get("caveat", ""), 0.7, 4.5, 11.9, 1.2, 11, MUTED)
    else:
        stat(sl, 0.7, 2.4, str(cov.get("total_circuits", "-")), "circuits modelled")
        stat(sl, 3.8, 2.4, str(cov.get("priced_circuits", "-")), "circuits priced",
             "Unpriced scope is excluded from the total, never estimated at a "
             "neighbouring rate.")
        stat(sl, 6.9, 2.4, f"{float(cov.get('effective_coverage_pct', 0)):.0%}",
             "effective coverage", cov.get("coverage_basis", "").replace("_", " ").lower())
        stat(sl, 10.0, 2.4, str(len(cov.get("unpriced_countries") or [])),
             "unpriced countries",
             ", ".join(cov.get("unpriced_countries") or []) or "none")

    # ------------------------------------------------------------- 4 scenarios
    if scen:
        sl = slide("Scenarios", "Savings by scenario",
                   "Scenarios differ in which levers they include, not in how "
                   "optimistic they are about the same levers.")
        rows = [["Scenario", "Current", "Target (base)", "Savings low",
                 "Savings base", "Savings high"]]
        for code in sorted(scen):
            v = scen[code]
            g = v["gross_run_rate_savings"]
            rows.append([code,
                         _money((snap.get("current_tco") or {}).get("total"), currency),
                         _money((v.get("target_tco") or {}).get("base"), currency),
                         _money(g["low"], currency), _money(g["base"], currency),
                         _money(g["high"], currency)])
        table(sl, rows, 0.7, 2.4, 11.9)

        # ---------------------------------------------------------- 5 levers
        sl = slide("Levers", f"What drives the saving - scenario {headline}",
                   "Each lever applies only to the cost layers it names, "
                   "compounding on what the previous lever left.")
        applied = scen[headline].get("levers_applied") or []
        rows = [["Lever", "Family", "Cost layers", "Saving (base)"]]
        for l in applied:
            rows.append([l.get("lever_id", ""), l.get("family", ""),
                         ", ".join(l.get("cost_layers") or []),
                         _money(l.get("saving_base"), currency)])
        if len(rows) > 1:
            table(sl, rows, 0.7, 2.4, 11.9, col_w=[2.6, 3.0, 2.3, 4.0])
        else:
            _text(sl, "No lever contributed a saving in this scenario.",
                  0.7, 2.4, 11.9, 0.4, 12, MUTED)

    # ----------------------------------------------------------- 6 confidence
    sl = slide("Confidence", "How much weight this estimate carries",
               "The band is derived from the run, not asserted. Where the "
               "evidence does not support a higher score, a ceiling is applied "
               "and named.")
    comps = conf.get("components") or {}
    rows = [["Component", "Score"]] + [[k.replace("_", " "), str(v)]
                                       for k, v in comps.items()]
    if len(rows) > 1:
        table(sl, rows, 0.7, 2.4, 5.6, col_w=[3.4, 2.2])
    ceilings = conf.get("ceilings_applied") or []
    _text(sl, "Ceilings applied", 7.0, 2.4, 5.6, 0.3, 12, INK, bold=True)
    _text(sl, "\n".join(f"- {c}" for c in ceilings) or
          "None - no ceiling was reached.", 7.0, 2.8, 5.6, 2.4, 11, MUTED)

    # ------------------------------------------------------------- 7 evidence
    sl = slide("Evidence", "What the 24-domain contract actually holds")
    summary = (disp.get("summary") or {}).get("counts") or {}
    rows = [["Disposition", "Domains"]] + [[k.replace("_", " "), str(v)]
                                           for k, v in summary.items() if v]
    if len(rows) > 1:
        table(sl, rows, 0.7, 2.4, 5.6, col_w=[3.4, 2.2])
    evidenced = [d for d in disp.get("dispositions") or []
                 if d.get("disposition") == "EVIDENCED_PUBLIC"]
    _text(sl, "Publicly evidenced domains", 7.0, 2.4, 5.6, 0.3, 12, INK, bold=True)
    _text(sl, "\n".join(f"- {d['domain_no']}. {d.get('domain_name')}"
                        for d in evidenced[:12]) or
          "None. Every input rests on a benchmark prior, an assertion or the "
          "simulated draw.", 7.0, 2.8, 5.6, 3.2, 11, MUTED)

    # --------------------------------------------------------- 8 what this isn't
    sl = slide("Limitations", "What this estimate is not",
               "Stated here rather than left to be discovered in the "
               "conversation this deck is presented in.")
    limits = [
        f"Stage V0. Confidence is ceilinged accordingly - there is no "
        f"contract, telemetry, serviceability or bid evidence behind it.",
        f"Simulated share {snap.get('simulated_share')}; asserted share "
        f"{snap.get('asserted_share')}. Both cap the score independently.",
    ]
    if method == "ANCHOR":
        limits += [
            "No site-level inventory was enumerated. The addressable share is "
            "a governed assumption and the single largest source of "
            "uncertainty here - not the saving rates.",
            "Circuit coverage reads 0% because nothing was enumerated, not "
            "because pricing failed. It is not comparable with a BUILD_UP run.",
        ]
    else:
        limits += [
            f"{len(cov.get('unpriced_pairs') or [])} (country, product, "
            f"bandwidth) pair(s) had no approved price and are excluded from "
            f"the total rather than estimated.",
        ]
    limits += [
        "One-time migration and transformation cost is not in the run-rate "
        "figure.",
        "Strategic fit, target architecture and design guardrails are not "
        "generated: the workbench holds no evidence for them, and this tool "
        "will not manufacture narrative a cost model cannot support.",
    ]
    _text(sl, "\n\n".join(f"- {t}" for t in limits), 0.7, 2.3, 11.9, 4.4, 12, INK)

    out = pathlib.Path(args.out)
    prs.save(str(out))
    print(f"wrote {out} ({len(prs.slides._sldIdLst)} slides)")
    print("Narrative slides - strategic fit, architecture, guardrails - are "
          "deliberately not generated. See the limitations slide.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
